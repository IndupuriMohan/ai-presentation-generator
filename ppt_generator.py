import streamlit as st
import time
import openai
from pptx import Presentation
from dotenv import load_dotenv
import os

# Load environment variables from the .env file
load_dotenv()

# Fetch the OpenAI API key
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OpenAI API key is missing. Please add it to a .env file.")
openai.api_key = api_key

# Function to generate a presentation
def generate_presentation(topic, num_slides):
    if num_slides > 10:
        num_slides = 10  # Limit slides to a maximum of 10

    # Create a new PowerPoint presentation
    prs = Presentation()

    for _ in range(num_slides):
        # Generate a subtopic
        prompt_subtopic = f"Generate a subtopic related to {topic}."
        response_subtopic = openai.Completion.create(
            model="gpt-3.5-turbo",
            prompt=prompt_subtopic,
            max_tokens=30,
            n=1
        )
        subtopic = response_subtopic.choices[0].text.strip()

        # Generate content for the subtopic
        prompt_content = f"Write about {subtopic} in the context of {topic} in not more than 50 words or in bullet points."
        response_content = openai.Completion.create(
            model="gpt-3.5-turbo",
            prompt=prompt_content,
            max_tokens=70
        )
        generated_content = response_content.choices[0].text.strip()

        # Add a slide to the presentation
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = subtopic
        content = slide.placeholders[1]
        content.text = generated_content

        # Pause to avoid rate limits
        time.sleep(10)

    # Save the presentation
    file_name = f"{topic}_presentation.pptx"
    prs.save(file_name)

    return file_name

# Streamlit App
def main():
    st.title("AI Presentation Generator")

    # User input
    user_topic = st.text_input("Enter the topic for the presentation:")
    user_num_slides = st.number_input("Enter the number of slides:", min_value=1, max_value=10, step=1)

    if st.button("Generate Presentation"):
        if user_topic.strip():
            st.info("Generating presentation... This may take a while.")
            try:
                presentation_file = generate_presentation(user_topic, int(user_num_slides))
                st.success(f"Presentation '{presentation_file}' generated successfully!")

                # Download button for the generated file
                with open(presentation_file, "rb") as file:
                    st.download_button(
                        label="Download Presentation",
                        data=file,
                        file_name=presentation_file,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            except Exception as e:
                st.error(f"An error occurred: {str(e)}")
        else:
            st.error("Please enter a topic for the presentation.")

if __name__ == "__main__":
    main()
